#!/usr/bin/env python3
"""
TripFlow 企业差旅报销流程管理平台 — PPT 生成脚本
基于模板 template.pptx 的主题/字体/配色，生成 38 页答辩演示文稿。
用法：python scripts/generate_ppt.py
输出：docs/汇报文件夹/TripFlow-答辩PPT.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu, Cm
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_AUTO_SIZE
import copy
import os

# ══════════════════════════════════════════════════════════════
# 常量与配色（继承模版主题色 + 补充项目色）
# ══════════════════════════════════════════════════════════════
SLIDE_W = Emu(12192000)  # 13.33 in
SLIDE_H = Emu(6858000)   # 7.50 in

# 主题色
COLOR_PRIMARY   = RGBColor(0x2B, 0x57, 0x9A)  # 深蓝
COLOR_ACCENT    = RGBColor(0x44, 0x72, 0xC4)  # 亮蓝
COLOR_BG        = RGBColor(0xF0, 0xF4, 0xF8)  # 浅灰蓝背景
COLOR_CARD_BG   = RGBColor(0xFF, 0xFF, 0xFF)  # 白色
COLOR_INNOVATION= RGBColor(0xF5, 0xA6, 0x23)  # 橙黄
COLOR_SUCCESS   = RGBColor(0x27, 0xAE, 0x60)  # 绿色
COLOR_ERROR     = RGBColor(0xE7, 0x4C, 0x3C)  # 红色
COLOR_TEXT      = RGBColor(0x33, 0x33, 0x33)  # 深灰
COLOR_TEXT_SEC  = RGBColor(0x66, 0x66, 0x66)  # 中灰
COLOR_WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
COLOR_DARK2     = RGBColor(0x44, 0x54, 0x6A)  # 模版 dk2
COLOR_LT2       = RGBColor(0xE7, 0xE6, 0xE6)  # 模版 lt2

# 字体
FONT_CN = "Microsoft YaHei"
FONT_EN = "Arial"
FONT_CODE = "Consolas"

# 路径
TEMPLATE_PATH = os.path.join(os.path.dirname(__file__), "..", "docs", "ppt-master-main", "projects", "template.pptx")
OUTPUT_DIR = os.path.join(os.path.dirname(__file__), "..", "docs", "汇报文件夹")
OUTPUT_PATH = os.path.join(OUTPUT_DIR, "TripFlow-答辩PPT.pptx")

# ══════════════════════════════════════════════════════════════
# 工具函数
# ══════════════════════════════════════════════════════════════

def Inches_val(v):
    return Inches(v)

def add_shape(slide, shape_type, left, top, width, height, fill_color=None, line_color=None, line_width=None):
    """添加一个形状"""
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    shape.shadow.inherit = False
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        if line_width:
            shape.line.width = line_width
    else:
        shape.line.fill.background()
    return shape

def add_textbox(slide, left, top, width, height, text="", font_size=14, bold=False,
                color=COLOR_TEXT, alignment=PP_ALIGN.LEFT, font_name=FONT_CN,
                anchor=MSO_ANCHOR.TOP):
    """添加文本框"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    txBox.text_frame.word_wrap = True
    txBox.text_frame.auto_size = MSO_AUTO_SIZE.NONE
    p = txBox.text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    # Set anchor
    from pptx.oxml.ns import qn
    txBox.text_frame._txBody.bodyPr.set('anchor', 't' if anchor == MSO_ANCHOR.TOP else 'ctr')
    return txBox

def add_paragraph(text_frame, text="", font_size=14, bold=False, color=COLOR_TEXT,
                  alignment=PP_ALIGN.LEFT, font_name=FONT_CN, space_before=0, space_after=0,
                  level=0):
    """在已有 text_frame 中追加段落"""
    p = text_frame.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    p.level = level
    if space_before:
        p.space_before = Pt(space_before)
    if space_after:
        p.space_after = Pt(space_after)
    return p

def add_run(paragraph, text, font_size=14, bold=False, color=COLOR_TEXT, font_name=FONT_CN):
    """在段落中添加 run"""
    run = paragraph.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font_name
    return run

def add_page_number(slide, page_num):
    """右下角添加页码"""
    add_textbox(slide, Inches(12.3), Inches(7.0), Inches(0.8), Inches(0.4),
                str(page_num), font_size=10, color=COLOR_TEXT_SEC, alignment=PP_ALIGN.RIGHT)

def add_title_bar(slide, title, top=Inches(0.3), subtitle=None):
    """页面顶部标题栏"""
    # 标题背景条
    add_shape(slide, MSO_SHAPE.RECTANGLE,
              Inches(0), top, SLIDE_W, Inches(0.7),
              fill_color=COLOR_PRIMARY)
    # 标题文字
    add_textbox(slide, Inches(0.6), top + Inches(0.08), Inches(10), Inches(0.55),
                title, font_size=24, bold=True, color=COLOR_WHITE)
    # 分隔线
    add_shape(slide, MSO_SHAPE.RECTANGLE,
              Inches(0), top + Inches(0.7), SLIDE_W, Inches(0.04),
              fill_color=COLOR_ACCENT)
    return top + Inches(0.74)

def add_subtitle_title(slide, title, top=Inches(0.3)):
    """无背景条的标题（直接文字）"""
    add_textbox(slide, Inches(0.6), top, Inches(10), Inches(0.6),
                title, font_size=24, bold=True, color=COLOR_PRIMARY)
    # 下划线
    add_shape(slide, MSO_SHAPE.RECTANGLE,
              Inches(0.6), top + Inches(0.55), Inches(1.5), Inches(0.04),
              fill_color=COLOR_ACCENT)
    return top + Inches(0.65)

def add_card(slide, left, top, width, height, fill=COLOR_CARD_BG, border_color=None,
             shadow=False, corner_radius=None):
    """添加圆角卡片"""
    shape = add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height,
                      fill_color=fill, line_color=border_color,
                      line_width=Pt(1) if border_color else None)
    if shadow:
        from pptx.oxml.ns import qn
        spPr = shape._element.spPr
        effect_lst = spPr.makeelement(qn('a:effectLst'), {})
        outer_shdw = effect_lst.makeelement(qn('a:outerShdw'), {
            'blurRad': '76200', 'dist': '38100', 'dir': '5400000',
            'algn': 'bl', 'rotWithShape': '0'
        })
        srgb = outer_shdw.makeelement(qn('a:srgbClr'), {'val': '000000'})
        alpha = srgb.makeelement(qn('a:alpha'), {'val': '23000'})
        srgb.append(alpha)
        outer_shdw.append(srgb)
        effect_lst.append(outer_shdw)
        spPr.append(effect_lst)
    return shape

def add_bullet_list(slide, left, top, width, height, items, font_size=13,
                    color=COLOR_TEXT, prefix="▸", bold_prefix=False):
    """添加带前缀的列表"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    txBox.text_frame.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = txBox.text_frame.paragraphs[0]
        else:
            p = txBox.text_frame.add_paragraph()
        p.space_before = Pt(4)
        p.space_after = Pt(2)
        if prefix:
            run1 = p.add_run()
            run1.text = f"{prefix} "
            run1.font.size = Pt(font_size)
            run1.font.bold = bold_prefix
            run1.font.color.rgb = COLOR_ACCENT
            run1.font.name = FONT_CN
        run2 = p.add_run()
        run2.text = item
        run2.font.size = Pt(font_size)
        run2.font.color.rgb = color
        run2.font.name = FONT_CN
    return txBox

def add_table(slide, left, top, width, height, rows, cols, data, col_widths=None,
              header_bg=COLOR_PRIMARY, header_fg=COLOR_WHITE):
    """添加表格"""
    table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table
    # Set column widths
    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = w
    # Header row
    for j, cell_text in enumerate(data[0]):
        cell = table.cell(0, j)
        cell.text = ""
        p = cell.text_frame.paragraphs[0]
        run = p.add_run()
        run.text = str(cell_text)
        run.font.size = Pt(12)
        run.font.bold = True
        run.font.color.rgb = header_fg
        run.font.name = FONT_CN
        p.alignment = PP_ALIGN.CENTER
        # Header bg
        from pptx.oxml.ns import qn
        tcPr = cell._tc.get_or_add_tcPr()
        solidFill = tcPr.makeelement(qn('a:solidFill'), {})
        srgb = solidFill.makeelement(qn('a:srgbClr'), {'val': str(header_bg)})
        solidFill.append(srgb)
        tcPr.append(solidFill)
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    # Data rows
    for i in range(1, rows):
        for j in range(cols):
            cell = table.cell(i, j)
            cell.text = ""
            p = cell.text_frame.paragraphs[0]
            run = p.add_run()
            run.text = str(data[i][j])
            run.font.size = Pt(11)
            run.font.color.rgb = COLOR_TEXT
            run.font.name = FONT_CN
            p.alignment = PP_ALIGN.CENTER
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            # Alternating row color
            if i % 2 == 0:
                tcPr = cell._tc.get_or_add_tcPr()
                solidFill = tcPr.makeelement(qn('a:solidFill'), {})
                srgb = solidFill.makeelement(qn('a:srgbClr'), {'val': 'F5F7FA'})
                solidFill.append(srgb)
                tcPr.append(solidFill)
    return table_shape

def set_cell_bg(cell, color_hex):
    """设置单元格背景色"""
    from pptx.oxml.ns import qn
    tcPr = cell._tc.get_or_add_tcPr()
    # Remove existing fill
    for existing in tcPr.findall(qn('a:solidFill')):
        tcPr.remove(existing)
    solidFill = tcPr.makeelement(qn('a:solidFill'), {})
    srgb = solidFill.makeelement(qn('a:srgbClr'), {'val': color_hex})
    solidFill.append(srgb)
    tcPr.append(solidFill)

def add_code_block(slide, left, top, width, height, code, font_size=11):
    """添加代码块（深色背景）"""
    bg = add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height,
                   fill_color=RGBColor(0x1E, 0x1E, 0x2E))
    txBox = slide.shapes.add_textbox(left + Inches(0.2), top + Inches(0.15),
                                     width - Inches(0.4), height - Inches(0.3))
    txBox.text_frame.word_wrap = True
    lines = code.strip().split('\n')
    for i, line in enumerate(lines):
        if i == 0:
            p = txBox.text_frame.paragraphs[0]
        else:
            p = txBox.text_frame.add_paragraph()
        p.space_before = Pt(1)
        p.space_after = Pt(1)
        run = p.add_run()
        run.text = line
        run.font.size = Pt(font_size)
        run.font.color.rgb = RGBColor(0xA6, 0xE2, 0x2E)  # green on dark
        run.font.name = FONT_CODE
    return bg

def add_green_check_items(slide, left, top, width, items, font_size=12):
    """添加绿色对勾列表项"""
    txBox = slide.shapes.add_textbox(left, top, width, Inches(len(items) * 0.35))
    txBox.text_frame.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = txBox.text_frame.paragraphs[0]
        else:
            p = txBox.text_frame.add_paragraph()
        p.space_before = Pt(3)
        run1 = p.add_run()
        run1.text = "✅ "
        run1.font.size = Pt(font_size)
        run1.font.color.rgb = COLOR_SUCCESS
        run1.font.name = FONT_CN
        run2 = p.add_run()
        run2.text = item
        run2.font.size = Pt(font_size)
        run2.font.color.rgb = COLOR_TEXT
        run2.font.name = FONT_CN
    return txBox

def add_flow_arrow(slide, left, top, width=Inches(0.5), height=Inches(0.25)):
    """添加流程箭头"""
    shape = add_shape(slide, MSO_SHAPE.RIGHT_ARROW, left, top, width, height,
                      fill_color=COLOR_ACCENT)
    return shape

def add_box_with_text(slide, left, top, width, height, text, bg_color=COLOR_ACCENT,
                      text_color=COLOR_WHITE, font_size=12, bold=True):
    """添加带文字的方框"""
    shape = add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height,
                      fill_color=bg_color)
    shape.text_frame.word_wrap = True
    p = shape.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = text_color
    run.font.name = FONT_CN
    p.alignment = PP_ALIGN.CENTER
    from pptx.oxml.ns import qn
    shape.text_frame._txBody.bodyPr.set('anchor', 'ctr')
    return shape

def add_warning_box(slide, left, top, width, height, text, bg_color=COLOR_ERROR):
    """添加警示框"""
    shape = add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height,
                      fill_color=bg_color)
    shape.text_frame.word_wrap = True
    p = shape.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = Pt(11)
    run.font.color.rgb = COLOR_WHITE
    run.font.name = FONT_CN
    from pptx.oxml.ns import qn
    shape.text_frame._txBody.bodyPr.set('anchor', 'ctr')
    return shape

def add_innovation_badge(slide, left, top):
    """添加 💡 创新标识"""
    add_textbox(slide, left, top, Inches(0.5), Inches(0.4),
                "💡", font_size=18, color=COLOR_INNOVATION)

def add_status_circle(slide, left, top, size, color, text=""):
    """添加状态圆形"""
    shape = add_shape(slide, MSO_SHAPE.OVAL, left, top, size, size, fill_color=color)
    if text:
        shape.text_frame.paragraphs[0].text = text
        shape.text_frame.paragraphs[0].font.size = Pt(int(size / Inches(1) * 8))
        shape.text_frame.paragraphs[0].font.color.rgb = COLOR_WHITE
        shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        from pptx.oxml.ns import qn
        shape.text_frame._txBody.bodyPr.set('anchor', 'ctr')
    return shape

def add_timeline_dot(slide, left, top, label, sublabel, is_active=False):
    """时间轴节点"""
    color = COLOR_PRIMARY if is_active else COLOR_ACCENT
    dot_size = Inches(0.3)
    add_status_circle(slide, left, top, dot_size, color)
    add_textbox(slide, left - Inches(0.6), top + Inches(0.35), Inches(1.5), Inches(0.3),
                label, font_size=11, bold=True, color=COLOR_PRIMARY, alignment=PP_ALIGN.CENTER)
    if sublabel:
        add_textbox(slide, left - Inches(0.6), top + Inches(0.55), Inches(1.5), Inches(0.3),
                    sublabel, font_size=9, color=COLOR_TEXT_SEC, alignment=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════
# 主生成函数
# ══════════════════════════════════════════════════════════════

def generate_presentation():
    # 加载模版（继承主题、字体、配色）
    prs = Presentation(TEMPLATE_PATH)
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    # 获取空白布局
    blank_layout = prs.slide_layouts[6]  # Blank
    title_layout = prs.slide_layouts[0]  # Title Slide
    title_content_layout = prs.slide_layouts[1]  # Title and Content

    # 删除模版中的所有现有幻灯片（保留布局和主题）
    while len(prs.slides) > 0:
        rId = prs.slides._sldIdLst[0].get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id')
        if rId is None:
            from pptx.oxml.ns import qn
            rId = prs.slides._sldIdLst[0].get(qn('r:id'))
        prs.part.drop_rel(rId)
        prs.slides._sldIdLst.remove(prs.slides._sldIdLst[0])

    page_num = 0

    # ════════════════════════════════════════════════════════
    # P1 — 封面页
    # ════════════════════════════════════════════════════════
    page_num += 1
    slide = prs.slides.add_slide(blank_layout)
    # 深蓝渐变背景
    bg = add_shape(slide, MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H,
                   fill_color=COLOR_PRIMARY)
    # 底部装饰条
    add_shape(slide, MSO_SHAPE.RECTANGLE, 0, Inches(6.6), SLIDE_W, Inches(0.9),
              fill_color=RGBColor(0x1A, 0x3A, 0x5C))
    # Logo/图标区域
    add_textbox(slide, Inches(4.5), Inches(1.0), Inches(4.3), Inches(1.0),
                "✈️ 💼", font_size=48, color=COLOR_WHITE, alignment=PP_ALIGN.CENTER)
    # 主标题
    add_textbox(slide, Inches(1.5), Inches(2.2), Inches(10.3), Inches(1.0),
                "企业差旅报销流程管理平台", font_size=36, bold=True,
                color=COLOR_WHITE, alignment=PP_ALIGN.CENTER)
    # 英文副标题
    add_textbox(slide, Inches(1.5), Inches(3.2), Inches(10.3), Inches(0.6),
                "Enterprise Travel Reimbursement System",
                font_size=16, color=RGBColor(0x8E, 0xB8, 0xE5), alignment=PP_ALIGN.CENTER)
    # 分隔线
    add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(4.5), Inches(4.0),
              Inches(4.3), Inches(0.02), fill_color=RGBColor(0x8E, 0xB8, 0xE5))
    # 团队信息
    add_textbox(slide, Inches(1.5), Inches(4.3), Inches(10.3), Inches(0.4),
                "团队：TripFlow", font_size=14, color=COLOR_WHITE, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(1.5), Inches(4.7), Inches(10.3), Inches(0.4),
                "讲者：吴汉东（组长）", font_size=14, color=COLOR_WHITE, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(1.5), Inches(5.1), Inches(10.3), Inches(0.4),
                "2026 年 6 月", font_size=13, color=RGBColor(0xAA, 0xCC, 0xEE),
                alignment=PP_ALIGN.CENTER)

    # ════════════════════════════════════════════════════════
    # P2 — 业务痛点
    # ════════════════════════════════════════════════════════
    page_num += 1
    slide = prs.slides.add_slide(blank_layout)
    content_top = add_title_bar(slide, "传统报销的四大痛点")
    add_page_number(slide, page_num)

    pain_points = [
        ("1", "审批周期长", "纸质单据层层签字，流转 3-5 天", "📋"),
        ("2", "真实性难验证", "缺乏系统化记录，无法追溯", "🔍"),
        ("3", "费用易出错", "手工计算容易算错，对账困难", "💰"),
        ("4", "反复沟通", "格式不对退回修改，效率低下", "🔄"),
    ]
    for i, (num, title, desc, icon) in enumerate(pain_points):
        y = content_top + Inches(0.4 + i * 1.05)
        # 编号圆形
        add_status_circle(slide, Inches(0.8), y, Inches(0.45), COLOR_ERROR, num)
        # 标题
        add_textbox(slide, Inches(1.5), y - Inches(0.02), Inches(4), Inches(0.35),
                    title, font_size=16, bold=True, color=COLOR_TEXT)
        # 描述
        add_textbox(slide, Inches(1.5), y + Inches(0.35), Inches(4), Inches(0.35),
                    desc, font_size=12, color=COLOR_TEXT_SEC)

    # 右侧图标区域
    right_x = Inches(7.0)
    for i, (num, title, desc, icon) in enumerate(pain_points):
        y = content_top + Inches(0.4 + i * 1.05)
        card = add_card(slide, right_x, y, Inches(5.0), Inches(0.85), fill=RGBColor(0xFD, 0xF2, 0xF2),
                       border_color=RGBColor(0xFA, 0xD7, 0xD7))
        add_textbox(slide, right_x + Inches(0.2), y + Inches(0.15), Inches(0.6), Inches(0.5),
                    icon, font_size=24, color=COLOR_ERROR)
        add_textbox(slide, right_x + Inches(0.9), y + Inches(0.2), Inches(3.8), Inches(0.4),
                    f"{title} — {desc}", font_size=12, color=COLOR_TEXT)

    # 底部总结
    add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0), Inches(6.8), SLIDE_W, Inches(0.04),
              fill_color=COLOR_ACCENT)
    add_textbox(slide, Inches(0.6), Inches(6.55), Inches(10), Inches(0.4),
                "→ TripFlow 为解决这些问题而生", font_size=16, bold=True, color=COLOR_PRIMARY)

    # ════════════════════════════════════════════════════════
    # P3 — 技术栈总览
    # ════════════════════════════════════════════════════════
    page_num += 1
    slide = prs.slides.add_slide(blank_layout)
    content_top = add_title_bar(slide, "技术架构")
    add_page_number(slide, page_num)

    # 表格
    table_data = [
        ["层级", "技术选型"],
        ["前端", "Vue 3 + TypeScript + Element Plus + Pinia"],
        ["后端", "Spring Boot 3.5 + MyBatis-Plus + JPA"],
        ["数据库", "MySQL 8（11 张表）"],
        ["缓存", "Redis 8（支持自动降级）"],
    ]
    add_table(slide, Inches(1.5), content_top + Inches(0.5), Inches(10), Inches(2.8),
              5, 2, table_data, col_widths=[Inches(2), Inches(8)])

    # 技术 Logo 区域（文字模拟）
    logos = [("Vue.js", COLOR_SUCCESS), ("Spring Boot", COLOR_SUCCESS),
             ("MySQL", COLOR_ACCENT), ("Redis", COLOR_ERROR)]
    for i, (name, color) in enumerate(logos):
        x = Inches(2.5 + i * 2.3)
        add_box_with_text(slide, x, content_top + Inches(3.8), Inches(1.8), Inches(0.6),
                         name, bg_color=color, font_size=13)

    # ════════════════════════════════════════════════════════
    # P4 — 核心功能模块
    # ════════════════════════════════════════════════════════
    page_num += 1
    slide = prs.slides.add_slide(blank_layout)
    content_top = add_title_bar(slide, "核心功能模块")
    add_page_number(slide, page_num)

    modules = [
        ("📋", "报销单列表", "多条件搜索 · 分页浏览"),
        ("📝", "报销单详情", "8 个折叠区块 · 完整流程"),
        ("📅", "补助日历", "逐天选择 · 精确计算"),
        ("💰", "费用分摊", "按项目分配 · 一键均摊"),
        ("📊", "工作台", "看板管理 · 待办/进行中/已完成"),
    ]
    card_w = Inches(3.6)
    card_h = Inches(1.8)
    gap = Inches(0.4)
    # Row 1: 3 cards
    for i in range(3):
        x = Inches(0.8 + i * (3.6 + 0.3))
        y = content_top + Inches(0.4)
        add_card(slide, x, y, card_w, card_h, shadow=True, border_color=COLOR_LT2)
        add_textbox(slide, x, y + Inches(0.15), card_w, Inches(0.6),
                    modules[i][0], font_size=32, alignment=PP_ALIGN.CENTER)
        add_textbox(slide, x + Inches(0.2), y + Inches(0.75), card_w - Inches(0.4), Inches(0.4),
                    modules[i][1], font_size=15, bold=True, color=COLOR_PRIMARY,
                    alignment=PP_ALIGN.CENTER)
        add_textbox(slide, x + Inches(0.2), y + Inches(1.2), card_w - Inches(0.4), Inches(0.4),
                    modules[i][2], font_size=11, color=COLOR_TEXT_SEC,
                    alignment=PP_ALIGN.CENTER)
    # Row 2: 2 cards (centered)
    for i in range(2):
        x = Inches(2.9 + i * (3.6 + 0.3))
        y = content_top + Inches(2.6)
        idx = i + 3
        add_card(slide, x, y, card_w, card_h, shadow=True, border_color=COLOR_LT2)
        add_textbox(slide, x, y + Inches(0.15), card_w, Inches(0.6),
                    modules[idx][0], font_size=32, alignment=PP_ALIGN.CENTER)
        add_textbox(slide, x + Inches(0.2), y + Inches(0.75), card_w - Inches(0.4), Inches(0.4),
                    modules[idx][1], font_size=15, bold=True, color=COLOR_PRIMARY,
                    alignment=PP_ALIGN.CENTER)
        add_textbox(slide, x + Inches(0.2), y + Inches(1.2), card_w - Inches(0.4), Inches(0.4),
                    modules[idx][2], font_size=11, color=COLOR_TEXT_SEC,
                    alignment=PP_ALIGN.CENTER)

    add_textbox(slide, Inches(0.6), Inches(6.6), Inches(8), Inches(0.3),
                "💬 此页后切换浏览器进行实际演示", font_size=10, color=COLOR_TEXT_SEC)

    # ════════════════════════════════════════════════════════
    # P5 — 团队分工
    # ════════════════════════════════════════════════════════
    page_num += 1
    slide = prs.slides.add_slide(blank_layout)
    content_top = add_title_bar(slide, "团队分工")
    add_page_number(slide, page_num)

    members = [
        ("吴汉东", "（组长）", ["前端全部", "后端基础设施", "统筹", "文档"], COLOR_PRIMARY),
        ("杨俊杰", "", ["报销单核心后端", "功能测试"], COLOR_ACCENT),
        ("吴云龙", "", ["补助日历", "数据库设计", "测试"], COLOR_SUCCESS),
    ]
    card_w = Inches(3.6)
    for i, (name, role, duties, color) in enumerate(members):
        x = Inches(0.8 + i * 4.1)
        y = content_top + Inches(0.4)
        # 卡片背景
        add_card(slide, x, y, card_w, Inches(4.2), shadow=True, border_color=COLOR_LT2)
        # 顶部色条
        add_shape(slide, MSO_SHAPE.RECTANGLE, x, y, card_w, Inches(0.06), fill_color=color)
        # 头像占位
        avatar_size = Inches(1.0)
        add_status_circle(slide, x + (card_w - avatar_size) / 2, y + Inches(0.3),
                         avatar_size, color, name[0])
        # 姓名
        add_textbox(slide, x, y + Inches(1.5), card_w, Inches(0.4),
                    name, font_size=20, bold=True, color=COLOR_TEXT, alignment=PP_ALIGN.CENTER)
        if role:
            add_textbox(slide, x, y + Inches(1.9), card_w, Inches(0.3),
                        role, font_size=12, color=COLOR_TEXT_SEC, alignment=PP_ALIGN.CENTER)
        # 职责列表
        add_bullet_list(slide, x + Inches(0.4), y + Inches(2.3), card_w - Inches(0.8), Inches(1.5),
                       duties, font_size=13, prefix="▸", bold_prefix=True)

    # 底部协作说明
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(6.6),
              Inches(12), Inches(0.6), fill_color=RGBColor(0xF5, 0xF7, 0xFA))
    add_textbox(slide, Inches(1.0), Inches(6.65), Inches(11.2), Inches(0.5),
                "协作方式：Git 分支管理 + 统一文档规范  |  模块纵向拆分，一人包到底",
                font_size=12, color=COLOR_TEXT_SEC, alignment=PP_ALIGN.CENTER)

    # ════════════════════════════════════════════════════════
    # P6 — 技术亮点
    # ════════════════════════════════════════════════════════
    page_num += 1
    slide = prs.slides.add_slide(blank_layout)
    content_top = add_title_bar(slide, "技术亮点")
    add_page_number(slide, page_num)

    highlights = [
        ("💡 银行家舍入", "四舍六入五成双", "避免系统性偏差，保证财务数据精确"),
        ("💡 乐观锁", "@Version 注解", "防并发修改冲突，版本冲突自动检测"),
        ("💡 Redis 缓存+降级", "三级 TTL 缓存", "双重删除策略，自动降级到内存缓存"),
        ("💡 自定义限流", "@RateLimit 注解", "Redis INCR 实现，自研非第三方框架"),
    ]
    card_w = Inches(5.8)
    card_h = Inches(2.2)
    for i in range(2):
        for j in range(2):
            idx = i * 2 + j
            x = Inches(0.6 + j * 6.2)
            y = content_top + Inches(0.3 + i * 2.4)
            # 卡片
            add_card(slide, x, y, card_w, card_h, shadow=True)
            # 左侧强调色条
            add_shape(slide, MSO_SHAPE.RECTANGLE, x, y, Inches(0.06), card_h,
                     fill_color=COLOR_INNOVATION)
            # 标题
            add_textbox(slide, x + Inches(0.3), y + Inches(0.15), card_w - Inches(0.5), Inches(0.4),
                        highlights[idx][0], font_size=15, bold=True, color=COLOR_INNOVATION)
            # 副标题
            add_textbox(slide, x + Inches(0.3), y + Inches(0.6), card_w - Inches(0.5), Inches(0.35),
                        highlights[idx][1], font_size=13, bold=True, color=COLOR_TEXT)
            # 描述
            add_textbox(slide, x + Inches(0.3), y + Inches(1.0), card_w - Inches(0.5), Inches(0.35),
                        highlights[idx][2], font_size=12, color=COLOR_TEXT_SEC)
            # 小示意条
            add_shape(slide, MSO_SHAPE.RECTANGLE, x + Inches(0.3), y + Inches(1.5),
                     Inches(3.0), Inches(0.04), fill_color=COLOR_LT2)

    # ════════════════════════════════════════════════════════
    # P7 — 文档交付物
    # ════════════════════════════════════════════════════════
    page_num += 1
    slide = prs.slides.add_slide(blank_layout)
    content_top = add_title_bar(slide, "项目文档体系")
    add_page_number(slide, page_num)

    # 左侧文档列表
    docs_list = [
        ("📄", "概要设计", "600+ 行"),
        ("📄", "详细设计", "模块 + 时序图"),
        ("📄", "数据库表结构", "11 表 + ER 图"),
        ("📄", "接口文档", "Swagger 生成"),
        ("📄", "自测报告", "脚本批量生成"),
    ]
    for i, (icon, name, desc) in enumerate(docs_list):
        y = content_top + Inches(0.3 + i * 0.9)
        add_textbox(slide, Inches(0.8), y, Inches(0.5), Inches(0.4),
                    icon, font_size=22)
        add_textbox(slide, Inches(1.4), y, Inches(2.5), Inches(0.35),
                    name, font_size=16, bold=True, color=COLOR_TEXT)
        add_textbox(slide, Inches(1.4), y + Inches(0.35), Inches(3.5), Inches(0.3),
                    desc, font_size=11, color=COLOR_TEXT_SEC)

    # 右侧目录截图占位
    right_x = Inches(6.0)
    add_card(slide, right_x, content_top + Inches(0.3), Inches(6.5), Inches(4.5),
            border_color=COLOR_LT2, shadow=True)
    files = ["1-API接口文档-v2.1.docx", "2-表结构定义-v1.1.docx",
             "3-开发自测报告-v3.0.docx", "4-开发项目WBS-v2.0.xlsx",
             "概要设计.md", "详细设计.docx"]
    for i, f in enumerate(files):
        y = content_top + Inches(0.6 + i * 0.6)
        add_textbox(slide, right_x + Inches(0.4), y, Inches(5.5), Inches(0.35),
                    f"📄 {f}", font_size=12, color=COLOR_TEXT)

    # ════════════════════════════════════════════════════════
    # P8 — 团队汇报过渡页
    # ════════════════════════════════════════════════════════
    page_num += 1
    slide = prs.slides.add_slide(blank_layout)
    add_page_number(slide, page_num)
    # 背景
    add_shape(slide, MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H, fill_color=COLOR_BG)
    # 主文字
    add_textbox(slide, Inches(1), Inches(2.0), Inches(11.3), Inches(0.8),
                "团队整体汇报完成", font_size=32, bold=True, color=COLOR_PRIMARY,
                alignment=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(1), Inches(2.9), Inches(11.3), Inches(0.6),
                "接下来进入个人专项汇报", font_size=18, color=COLOR_TEXT_SEC,
                alignment=PP_ALIGN.CENTER)
    # 三个人物指示卡片
    person_items = [
        ("吴汉东", "前端+基础", COLOR_PRIMARY, True),
        ("杨俊杰", "核心后端", COLOR_ACCENT, False),
        ("吴云龙", "日历+DB", COLOR_SUCCESS, False),
    ]
    for i, (name, desc, color, highlight) in enumerate(person_items):
        x = Inches(2.5 + i * 3.0)
        y = Inches(4.0)
        sz = Inches(1.3) if highlight else Inches(1.1)
        ox = (sz - Inches(1.1)) / 2
        add_card(slide, x, y - ox, sz + Inches(0.8), sz + Inches(0.8),
                fill=COLOR_CARD_BG, border_color=color, shadow=highlight)
        add_status_circle(slide, x + Inches(0.25), y + Inches(0.15 - ox), Inches(0.5),
                         color, name[0])
        add_textbox(slide, x, y + Inches(0.7 - ox), Inches(0.8 + sz), Inches(0.3),
                    name, font_size=12 if not highlight else 14, bold=True, color=color,
                    alignment=PP_ALIGN.CENTER)
        add_textbox(slide, x, y + Inches(1.0 - ox), Inches(0.8 + sz), Inches(0.25),
                    desc, font_size=10, color=COLOR_TEXT_SEC, alignment=PP_ALIGN.CENTER)
        # 箭头
        if i < 2:
            add_flow_arrow(slide, x + sz + Inches(0.6), y + Inches(0.3),
                          Inches(0.4), Inches(0.2))

    # 底部指示
    add_textbox(slide, Inches(1), Inches(6.0), Inches(11.3), Inches(0.5),
                "↓ 现在进入吴汉东的个人专项 ↓", font_size=14, color=COLOR_PRIMARY,
                alignment=PP_ALIGN.CENTER, bold=True)

    # ════════════════════════════════════════════════════════
    # Part B — 吴汉东个人专项
    # ════════════════════════════════════════════════════════

    # P9 — 吴汉东个人封面
    page_num += 1
    slide = prs.slides.add_slide(blank_layout)
    add_shape(slide, MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H, fill_color=COLOR_BG)
    add_page_number(slide, page_num)
    add_textbox(slide, Inches(1), Inches(1.5), Inches(11.3), Inches(0.8),
                "吴汉东 个人专项汇报", font_size=32, bold=True, color=COLOR_PRIMARY,
                alignment=PP_ALIGN.CENTER)

    modules_whd = [
        ("🎨", "前端开发", "Vue3 全栈页面", COLOR_PRIMARY),
        ("⚙️", "后端基础设施", "缓存/限流/异常", COLOR_ACCENT),
        ("📄", "文档撰写", "概要+详细设计", COLOR_SUCCESS),
    ]
    for i, (icon, title, desc, color) in enumerate(modules_whd):
        x = Inches(1.5 + i * 3.8)
        y = Inches(3.0)
        add_card(slide, x, y, Inches(3.2), Inches(2.5), shadow=True, border_color=color)
        add_textbox(slide, x, y + Inches(0.3), Inches(3.2), Inches(0.6),
                    icon, font_size=36, alignment=PP_ALIGN.CENTER)
        add_textbox(slide, x, y + Inches(1.0), Inches(3.2), Inches(0.4),
                    title, font_size=16, bold=True, color=color, alignment=PP_ALIGN.CENTER)
        add_textbox(slide, x, y + Inches(1.5), Inches(3.2), Inches(0.4),
                    desc, font_size=12, color=COLOR_TEXT_SEC, alignment=PP_ALIGN.CENTER)

    # P10 — 前端项目架构
    page_num += 1
    slide = prs.slides.add_slide(blank_layout)
    content_top = add_title_bar(slide, "前端技术架构")
    add_page_number(slide, page_num)

    # 左侧目录树
    tree_lines = [
        "tripflow-web/src/",
        "├── api/          → HTTP 请求层",
        "├── stores/       → 状态管理层",
        "├── composables/  → 可复用逻辑",
        "├── types/        → 类型定义",
        "├── utils/        → 工具函数",
        "├── views/        → 页面组件",
        "├── components/   → 子组件",
        "└── router/       → 路由配置",
    ]
    add_code_block(slide, Inches(0.6), content_top + Inches(0.3), Inches(5.5), Inches(4.2),
                   "\n".join(tree_lines), font_size=12)

    # 右侧分层说明
    right_x = Inches(6.8)
    layers = [
        ("API 层", "HTTP 请求封装", COLOR_PRIMARY),
        ("Store 层", "Pinia 状态管理", COLOR_ACCENT),
        ("Composable 层", "可复用逻辑", COLOR_SUCCESS),
        ("View/Component", "页面和组件", COLOR_INNOVATION),
    ]
    for i, (name, desc, color) in enumerate(layers):
        y = content_top + Inches(0.4 + i * 1.05)
        add_box_with_text(slide, right_x, y, Inches(2.2), Inches(0.6),
                         name, bg_color=color, font_size=12)
        add_textbox(slide, right_x + Inches(2.5), y + Inches(0.1), Inches(3.0), Inches(0.4),
                    desc, font_size=11, color=COLOR_TEXT_SEC)
        if i < 3:
            add_textbox(slide, right_x + Inches(0.7), y + Inches(0.65), Inches(1), Inches(0.3),
                        "↓", font_size=14, color=COLOR_TEXT_SEC, alignment=PP_ALIGN.CENTER)

    # 底部指标
    metrics = ["✅ 4000+ 行 TypeScript", "✅ TypeScript 类型全覆盖", "✅ 单文件不超过 200 行"]
    for i, m in enumerate(metrics):
        add_textbox(slide, Inches(0.6 + i * 4.0), Inches(6.6), Inches(3.8), Inches(0.3),
                    m, font_size=11, color=COLOR_SUCCESS, bold=True)

    # P11 — 报销单列表页展示
    page_num += 1
    slide = prs.slides.add_slide(blank_layout)
    content_top = add_title_bar(slide, "报销单列表页")
    add_page_number(slide, page_num)

    # 系统截图占位
    img_w = Inches(8.5)
    img_h = Inches(4.8)
    img_x = (SLIDE_W - img_w) / 2
    add_card(slide, img_x, content_top + Inches(0.2), img_w, img_h,
            border_color=COLOR_LT2, shadow=True)
    # 截图区域说明
    add_textbox(slide, img_x + Inches(0.5), content_top + Inches(0.5), Inches(7.5), Inches(3.5),
                "📷 系统截图 — 报销单列表页\n\n"
                "顶部：7 个搜索条件\n"
                "  （单据编号、出差事由、报销人、部门、公司、业务类型、状态）\n\n"
                "中部：数据表格，分页显示\n\n"
                "右侧操作列：详情 / 编辑 / 复制 / 提交 / 作废",
                font_size=13, color=COLOR_TEXT_SEC, alignment=PP_ALIGN.LEFT)
    # 技术要点
    add_textbox(slide, Inches(0.6), Inches(6.6), Inches(11), Inches(0.3),
                "reactive 统一管理 · URL 同步分页 · 下拉菜单操作",
                font_size=11, color=COLOR_TEXT_SEC, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.6), Inches(6.9), Inches(11), Inches(0.3),
                "💬 切换浏览器演示实际列表页",
                font_size=10, color=COLOR_ACCENT, alignment=PP_ALIGN.CENTER)

    # P12 — 报销单详情页展示
    page_num += 1
    slide = prs.slides.add_slide(blank_layout)
    content_top = add_title_bar(slide, "报销单详情页（8 个可折叠区块）")
    add_page_number(slide, page_num)

    # 左侧截图
    img_w = Inches(6.5)
    img_h = Inches(5.0)
    add_card(slide, Inches(0.6), content_top + Inches(0.2), img_w, img_h,
            border_color=COLOR_LT2, shadow=True)
    add_textbox(slide, Inches(1.0), content_top + Inches(0.5), Inches(5.5), Inches(4.0),
                "📷 系统截图 — 详情页全貌\n（展示所有展开的折叠区块）",
                font_size=13, color=COLOR_TEXT_SEC, alignment=PP_ALIGN.CENTER)

    # 右侧 8 个区块
    blocks = ["① 表头", "② 基本信息", "③ 行程管理", "④ 补助信息",
              "⑤ 费用合计", "⑥ 费用分摊", "⑦ 备注", "⑧ 操作栏"]
    right_x = Inches(7.5)
    for i, b in enumerate(blocks):
        y = content_top + Inches(0.3 + i * 0.55)
        add_box_with_text(slide, right_x, y, Inches(4.5), Inches(0.42),
                         b, bg_color=COLOR_ACCENT if i < 4 else COLOR_PRIMARY,
                         font_size=11)

    # 底部特性
    features = ["新建/查看/编辑模式复用同一组件", "防重复提交（按钮 loading + 请求锁）",
                "提交成功动画进度条"]
    add_green_check_items(slide, Inches(0.6), Inches(6.4), Inches(12), features, font_size=11)

    # P13 — 补助日历弹窗展示
    page_num += 1
    slide = prs.slides.add_slide(blank_layout)
    content_top = add_title_bar(slide, "补助日历弹窗")
    add_page_number(slide, page_num)

    # 截图
    img_w = Inches(8.0)
    img_h = Inches(4.0)
    img_x = (SLIDE_W - img_w) / 2
    add_card(slide, img_x, content_top + Inches(0.2), img_w, img_h,
            border_color=COLOR_LT2, shadow=True)
    add_textbox(slide, img_x + Inches(0.5), content_top + Inches(0.5), Inches(7.0), Inches(3.0),
                "📷 系统截图 — 补助日历弹窗\n\n"
                "行（行程）× 列（日期）矩阵\n"
                "每天可勾选：餐补 / 交通补 / 通讯补",
                font_size=13, color=COLOR_TEXT_SEC, alignment=PP_ALIGN.CENTER)

    # 底部功能标签
    func_labels = [("📅", "逐天选择"), ("🖱️", "批量勾选"), ("💰", "实时计算")]
    for i, (icon, label) in enumerate(func_labels):
        x = Inches(2.5 + i * 3.0)
        y = content_top + Inches(4.5)
        add_card(slide, x, y, Inches(2.4), Inches(1.0), shadow=True, border_color=COLOR_ACCENT)
        add_textbox(slide, x, y + Inches(0.05), Inches(2.4), Inches(0.5),
                    icon, font_size=24, alignment=PP_ALIGN.CENTER)
        add_textbox(slide, x, y + Inches(0.55), Inches(2.4), Inches(0.35),
                    label, font_size=13, bold=True, color=COLOR_PRIMARY,
                    alignment=PP_ALIGN.CENTER)

    add_textbox(slide, Inches(0.6), Inches(6.8), Inches(11), Inches(0.3),
                "💬 切换浏览器打开日历弹窗演示",
                font_size=10, color=COLOR_ACCENT, alignment=PP_ALIGN.CENTER)

    # P14 — Redis 缓存架构
    page_num += 1
    slide = prs.slides.add_slide(blank_layout)
    content_top = add_title_bar(slide, "Redis 缓存架构")
    add_innovation_badge(slide, Inches(5.5), content_top - Inches(0.55))
    add_page_number(slide, page_num)

    # 左侧架构流程图
    flow_items = [("Client", COLOR_DARK2), ("Spring Boot", COLOR_PRIMARY),
                  ("Redis Cache", COLOR_ERROR), ("内存缓存 (fallback)", COLOR_ACCENT)]
    flow_x = Inches(0.6)
    for i, (name, color) in enumerate(flow_items):
        y = content_top + Inches(0.3 + i * 1.0)
        add_box_with_text(slide, flow_x, y, Inches(3.5), Inches(0.6),
                         name, bg_color=color, font_size=13)
        if i < 3:
            line_style = "dashed" if i == 2 else "solid"
            add_textbox(slide, flow_x + Inches(1.3), y + Inches(0.6), Inches(1), Inches(0.35),
                        "↓" if i < 2 else "↓ (降级)", font_size=11, color=COLOR_TEXT_SEC,
                        alignment=PP_ALIGN.CENTER)

    # 右侧 — 三级 TTL
    right_x = Inches(5.0)
    add_textbox(slide, right_x, content_top + Inches(0.2), Inches(3.5), Inches(0.35),
                "三级缓存 TTL", font_size=15, bold=True, color=COLOR_PRIMARY)
    ttl_data = [
        ["数据类型", "TTL"],
        ["基础数据", "12 小时"],
        ["报销人", "1 小时"],
        ["报销单", "60 秒"],
    ]
    add_table(slide, right_x, content_top + Inches(0.6), Inches(3.5), Inches(2.0),
              4, 2, ttl_data, col_widths=[Inches(1.8), Inches(1.7)])

    # 双重删除策略
    add_textbox(slide, right_x + Inches(4.2), content_top + Inches(0.2), Inches(3.5), Inches(0.35),
                "双重删除策略", font_size=15, bold=True, color=COLOR_PRIMARY)
    del_items = ["① 更新前删除缓存", "② 提交后删除缓存", "③ 300ms 延迟删除"]
    add_bullet_list(slide, right_x + Inches(4.2), content_top + Inches(0.6),
                   Inches(3.5), Inches(1.5), del_items, font_size=12)

    # 降级说明
    add_warning_box(slide, right_x, content_top + Inches(3.0), Inches(7.5), Inches(0.5),
                   "⚠️ 自动降级：CacheErrorHandler 捕获异常 → 降级到内存缓存")

    # P15 — 自定义限流注解
    page_num += 1
    slide = prs.slides.add_slide(blank_layout)
    content_top = add_title_bar(slide, "@RateLimit 自定义限流注解")
    add_innovation_badge(slide, Inches(7.2), content_top - Inches(0.55))
    add_page_number(slide, page_num)

    # 左侧代码
    code = '''@RateLimit(
  key = "submit",
  maxRequests = 10,
  timeWindowSeconds = 60
)
public ApiResponse<Void>
  submitReimbursement(...) {
  // 业务逻辑
}'''
    add_code_block(slide, Inches(0.6), content_top + Inches(0.3), Inches(5.5), Inches(3.5),
                   code, font_size=12)

    # 右侧原理
    right_x = Inches(6.8)
    add_textbox(slide, right_x, content_top + Inches(0.2), Inches(3), Inches(0.35),
                "实现原理", font_size=15, bold=True, color=COLOR_PRIMARY)
    steps = [
        "1. Redis INCR 计数器 +1",
        "2. 超过阈值 → 返回 429",
        "3. EXPIRE 设置窗口过期时间",
    ]
    add_bullet_list(slide, right_x, content_top + Inches(0.6), Inches(5.5), Inches(1.5),
                   steps, font_size=13, prefix="", bold_prefix=True)

    # 底部特性
    add_green_check_items(slide, Inches(0.6), Inches(6.4), Inches(12),
                         ["基于 Spring AOP + HandlerInterceptor",
                          "自研注解，不依赖 Sentinel/Guava"],
                         font_size=12)

    # P16 — 统一异常处理
    page_num += 1
    slide = prs.slides.add_slide(blank_layout)
    content_top = add_title_bar(slide, "统一异常处理（wrap() 方法）")
    add_page_number(slide, page_num)

    exc_data = [
        ["异常类型", "状态码", "触发场景"],
        ["ReimbursementVersionConflictException", "409", "乐观锁冲突"],
        ['IllegalArgumentException（含“不存在”）', "404", "资源不存在"],
        ["IllegalArgumentException（其他）", "400", "参数校验失败"],
        ["DateTimeParseException", "400", "日期格式错误"],
    ]
    add_table(slide, Inches(0.6), content_top + Inches(0.3), Inches(12), Inches(2.5),
              5, 3, exc_data, col_widths=[Inches(5.5), Inches(1.5), Inches(3.0)])

    add_textbox(slide, Inches(0.6), content_top + Inches(3.2), Inches(12), Inches(0.4),
                "前端 Axios 拦截器 → 自动提取错误信息 → 展示给用户",
                font_size=13, color=COLOR_PRIMARY, bold=True)

    # 截图占位
    add_card(slide, Inches(3.0), content_top + Inches(3.8), Inches(6.0), Inches(1.5),
            border_color=COLOR_LT2)
    add_textbox(slide, Inches(3.5), content_top + Inches(4.0), Inches(5.0), Inches(1.0),
                "📷 前端错误提示弹窗截图",
                font_size=12, color=COLOR_TEXT_SEC, alignment=PP_ALIGN.CENTER)

    # P17 — 吴汉东小结
    page_num += 1
    slide = prs.slides.add_slide(blank_layout)
    content_top = add_title_bar(slide, "吴汉东 — 模块总结")
    add_page_number(slide, page_num)

    summary_items = [
        ("🎨 前端", ["7 个页面", "10 个组件", "TS 全覆盖", "4000+ 行"], COLOR_PRIMARY),
        ("⚙️ 后端", ["Redis 缓存架构", "限流注解", "异常处理"], COLOR_ACCENT),
        ("📄 文档", ["概要设计", "详细设计", "自测报告"], COLOR_SUCCESS),
    ]
    for i, (title, items, color) in enumerate(summary_items):
        x = Inches(0.6 + i * 4.2)
        y = content_top + Inches(0.4)
        add_card(slide, x, y, Inches(3.8), Inches(3.5), shadow=True, border_color=color)
        add_shape(slide, MSO_SHAPE.RECTANGLE, x, y, Inches(3.8), Inches(0.06),
                 fill_color=color)
        add_textbox(slide, x, y + Inches(0.2), Inches(3.8), Inches(0.4),
                    title, font_size=18, bold=True, color=color, alignment=PP_ALIGN.CENTER)
        add_bullet_list(slide, x + Inches(0.4), y + Inches(0.8), Inches(3.0), Inches(2.5),
                       items, font_size=14, prefix="▸")

    # P18 — 过渡页（→ 杨俊杰）
    page_num += 1
    slide = prs.slides.add_slide(blank_layout)
    add_shape(slide, MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H, fill_color=COLOR_BG)
    add_page_number(slide, page_num)
    add_textbox(slide, Inches(1), Inches(2.2), Inches(11.3), Inches(0.8),
                "吴汉东 专项汇报完成 ✓", font_size=28, bold=True, color=COLOR_PRIMARY,
                alignment=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(1), Inches(3.2), Inches(11.3), Inches(0.5),
                "接下来 ↓", font_size=18, color=COLOR_TEXT_SEC, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(1), Inches(4.0), Inches(11.3), Inches(0.6),
                "杨俊杰 — 报销单核心业务后端", font_size=22, bold=True, color=COLOR_ACCENT,
                alignment=PP_ALIGN.CENTER)

    # ════════════════════════════════════════════════════════
    # Part C — 杨俊杰个人专项
    # ════════════════════════════════════════════════════════

    # P19 — 杨俊杰个人封面
    page_num += 1
    slide = prs.slides.add_slide(blank_layout)
    add_shape(slide, MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H, fill_color=COLOR_BG)
    add_page_number(slide, page_num)
    add_textbox(slide, Inches(1), Inches(1.5), Inches(11.3), Inches(0.8),
                "杨俊杰 个人专项汇报", font_size=32, bold=True, color=COLOR_PRIMARY,
                alignment=PP_ALIGN.CENTER)

    yjj_modules = [
        "1. 报销单基本信息管理",
        "2. 补录行程管理",
        "3. 费用合计与分摊",
        "4. 报销单状态管理",
        "5. 功能/性能/安全测试",
    ]
    add_card(slide, Inches(3.0), Inches(2.8), Inches(7.0), Inches(3.0),
            shadow=True, border_color=COLOR_ACCENT)
    add_textbox(slide, Inches(3.5), Inches(2.9), Inches(6.0), Inches(0.4),
                "负责模块：", font_size=14, bold=True, color=COLOR_ACCENT)
    add_bullet_list(slide, Inches(3.5), Inches(3.4), Inches(6.0), Inches(2.3),
                   yjj_modules, font_size=15, prefix="▸", bold_prefix=True)

    # P20 — 报销单 CRUD 接口
    page_num += 1
    slide = prs.slides.add_slide(blank_layout)
    content_top = add_title_bar(slide, "报销单核心 CRUD — 10 个接口")
    add_page_number(slide, page_num)

    api_data = [
        ["方法", "路径", "功能"],
        ["GET", "/reimbursement", "分页列表（7 个筛选条件）"],
        ["GET", "/reimbursement/{id}", "详情（含所有子表数据）"],
        ["POST", "/reimbursement", "创建（默认草稿状态）"],
        ["PUT", "/reimbursement/{id}", "更新（需携带 version）"],
        ["DELETE", "/reimbursement/{id}", "删除（需携带 version）"],
        ["POST", ".../submit", "提交审批"],
        ["POST", ".../void", "作废"],
        ["PUT", ".../remark", "更新备注"],
        ["DELETE", ".../remark", "清除备注"],
        ["GET", ".../expense-summary", "费用合计"],
    ]
    add_table(slide, Inches(0.6), content_top + Inches(0.2), Inches(12), Inches(4.5),
              11, 3, api_data, col_widths=[Inches(1.5), Inches(5.0), Inches(4.5)])

    # 💡 标注
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(6.4),
              Inches(12), Inches(0.5), fill_color=RGBColor(0xFE, 0xF9, 0xE7))
    add_textbox(slide, Inches(1.0), Inches(6.43), Inches(11), Inches(0.4),
                "💡 content 字段用 JSON 存储主体信息 → 灵活性 + 查询性能兼顾",
                font_size=12, color=COLOR_INNOVATION, bold=True)

    # P21 — 数据表结构（JSON 设计）
    page_num += 1
    slide = prs.slides.add_slide(blank_layout)
    content_top = add_title_bar(slide, "数据存储设计")
    add_page_number(slide, page_num)

    # 左侧表结构树
    tree = """reimbursement (主表)
├── id, status, version
├── content (JSON)
├── reimburser_name
├── department_name, company_name
│
├── reimbursement_travel_record
├── reimbursement_allowance
│   └── reimbursement_allowance_calendar
└── reimbursement_cost_allocation"""
    add_code_block(slide, Inches(0.6), content_top + Inches(0.3), Inches(6.0), Inches(4.2),
                   tree, font_size=12)

    # 右侧 JSON 设计说明
    right_x = Inches(7.2)
    add_textbox(slide, right_x, content_top + Inches(0.3), Inches(5), Inches(0.4),
                "为什么用 JSON？", font_size=16, bold=True, color=COLOR_PRIMARY)
    json_reasons = [
        "✅ 灵活性：新字段无需 ALTER TABLE",
        "✅ 查询性能：主表查询不需 JOIN",
        "✅ 搜索字段冗余到独立字段",
        "    （reimburser_name, department_name 等）",
    ]
    add_bullet_list(slide, right_x, content_top + Inches(0.9), Inches(5.5), Inches(3.0),
                   json_reasons, font_size=13, prefix="", bold_prefix=False)

    # P22 — 行程重叠检测
    page_num += 1
    slide = prs.slides.add_slide(blank_layout)
    content_top = add_title_bar(slide, "分钟级行程重叠检测")
    add_innovation_badge(slide, Inches(6.5), content_top - Inches(0.55))
    add_page_number(slide, page_num)

    # 时间轴示意图
    # 行程 A
    add_box_with_text(slide, Inches(1.0), content_top + Inches(0.5),
                     Inches(6.0), Inches(0.6),
                     "行程 A: [startA ──────────── endA]", bg_color=COLOR_ACCENT, font_size=12)
    # 行程 B
    add_box_with_text(slide, Inches(3.0), content_top + Inches(1.3),
                     Inches(6.0), Inches(0.6),
                     "行程 B:     [startB ──── endB]", bg_color=COLOR_ERROR, font_size=12)

    # 公式
    add_textbox(slide, Inches(0.6), content_top + Inches(2.2), Inches(12), Inches(0.5),
                "重叠条件: startA ≤ endB  AND  startB ≤ endA",
                font_size=16, bold=True, color=COLOR_PRIMARY, alignment=PP_ALIGN.CENTER)

    # 对比示例
    # 合法
    add_card(slide, Inches(0.6), content_top + Inches(2.9), Inches(5.5), Inches(2.2),
            border_color=COLOR_SUCCESS)
    add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0.6), content_top + Inches(2.9),
              Inches(5.5), Inches(0.06), fill_color=COLOR_SUCCESS)
    add_textbox(slide, Inches(1.0), content_top + Inches(3.1), Inches(4.5), Inches(0.4),
                "✅ 合法", font_size=16, bold=True, color=COLOR_SUCCESS)
    add_textbox(slide, Inches(1.0), content_top + Inches(3.6), Inches(4.5), Inches(1.2),
                "A: 6/1 08:00 — 6/5 18:00\nB: 6/5 19:00 — 6/8 20:00\n→ 同日不同时，不冲突",
                font_size=12, color=COLOR_TEXT)

    # 冲突
    add_card(slide, Inches(6.8), content_top + Inches(2.9), Inches(5.5), Inches(2.2),
            border_color=COLOR_ERROR)
    add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(6.8), content_top + Inches(2.9),
              Inches(5.5), Inches(0.06), fill_color=COLOR_ERROR)
    add_textbox(slide, Inches(7.2), content_top + Inches(3.1), Inches(4.5), Inches(0.4),
                "❌ 冲突", font_size=16, bold=True, color=COLOR_ERROR)
    add_textbox(slide, Inches(7.2), content_top + Inches(3.6), Inches(4.5), Inches(1.2),
                "A: 6/1 — 6/5\nB: 6/3 — 6/8\n→ 时间重叠，返回 400 错误",
                font_size=12, color=COLOR_TEXT)

    # P23 — 费用合计与分摊
    page_num += 1
    slide = prs.slides.add_slide(blank_layout)
    content_top = add_title_bar(slide, "费用合计与分摊")
    add_page_number(slide, page_num)

    # 左侧 JSON
    json_code = '''{
  "totalAllowance": 2400.00,
  "totalMeal": 1500.00,
  "totalTransport": 600.00,
  "totalCommunication": 300.00,
  "grandTotal": 2400.00
}'''
    add_textbox(slide, Inches(0.6), content_top + Inches(0.2), Inches(5), Inches(0.35),
                "费用合计响应", font_size=15, bold=True, color=COLOR_PRIMARY)
    add_code_block(slide, Inches(0.6), content_top + Inches(0.6), Inches(5.5), Inches(2.8),
                   json_code, font_size=13)

    # 右侧分摊规则
    right_x = Inches(6.8)
    add_textbox(slide, right_x, content_top + Inches(0.2), Inches(5), Inches(0.35),
                "费用分摊规则", font_size=15, bold=True, color=COLOR_PRIMARY)
    rules = [
        "1. 首行锁定：第 1 行 = 总额 - 其他行之和",
        "2. 均摊功能：总额 ÷ N 行",
        "3. 比例上限：合计 = 100%",
        "4. 银行家舍入：四舍六入五成双",
    ]
    add_bullet_list(slide, right_x, content_top + Inches(0.7), Inches(5.5), Inches(2.0),
                   rules, font_size=13, prefix="", bold_prefix=True)

    # 底部💡
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(5.8),
              Inches(12), Inches(0.6), fill_color=RGBColor(0xFE, 0xF9, 0xE7))
    add_textbox(slide, Inches(1.0), Inches(5.85), Inches(11), Inches(0.5),
                "💡 首行锁定保证总额不错：1000 元分 3 项各 33% → 330 + 330 + 340 = 1000",
                font_size=12, color=COLOR_INNOVATION, bold=True)

    # P24 — 费用分摊接口
    page_num += 1
    slide = prs.slides.add_slide(blank_layout)
    content_top = add_title_bar(slide, "费用分摊 — 5 个接口")
    add_page_number(slide, page_num)

    alloc_data = [
        ["方法", "路径", "功能"],
        ["GET", ".../cost-allocations", "分摊列表"],
        ["POST", ".../cost-allocations", "新增分摊"],
        ["POST", ".../evenly-distribute", "一键均摊"],
        ["PUT", ".../cost-allocations/{key}", "更新分摊"],
        ["DELETE", ".../cost-allocations/{key}", "删除分摊"],
    ]
    add_table(slide, Inches(1.5), content_top + Inches(0.3), Inches(9.5), Inches(3.0),
              6, 3, alloc_data, col_widths=[Inches(1.5), Inches(4.5), Inches(3.0)])

    # 截图占位
    add_card(slide, Inches(2.0), content_top + Inches(3.5), Inches(8.5), Inches(1.8),
            border_color=COLOR_LT2)
    add_textbox(slide, Inches(3.0), content_top + Inches(3.8), Inches(6.5), Inches(1.0),
                "📷 费用分摊界面截图", font_size=13, color=COLOR_TEXT_SEC,
                alignment=PP_ALIGN.CENTER)

    # P25 — 状态管理与乐观锁
    page_num += 1
    slide = prs.slides.add_slide(blank_layout)
    content_top = add_title_bar(slide, "状态管理 + 乐观锁并发控制")
    add_page_number(slide, page_num)

    # 左侧状态流转
    states = [("草稿(0)", COLOR_ACCENT), ("已提交(1)", COLOR_PRIMARY), ("已作废(2)", COLOR_ERROR)]
    for i, (name, color) in enumerate(states):
        x = Inches(0.8 + i * 2.5)
        y = content_top + Inches(0.8)
        add_box_with_text(slide, x, y, Inches(1.8), Inches(0.6),
                         name, bg_color=color, font_size=13)
        if i < 2:
            add_textbox(slide, x + Inches(1.9), y + Inches(0.1), Inches(0.5), Inches(0.4),
                        "→", font_size=20, color=COLOR_TEXT_SEC)
            action = "提交" if i == 0 else "作废"
            add_textbox(slide, x + Inches(1.85), y - Inches(0.35), Inches(0.7), Inches(0.3),
                        action, font_size=10, color=COLOR_TEXT_SEC, alignment=PP_ALIGN.CENTER)

    # 右侧乐观锁
    right_x = Inches(7.5)
    add_textbox(slide, right_x, content_top + Inches(0.2), Inches(5), Inches(0.35),
                "乐观锁实现", font_size=15, bold=True, color=COLOR_PRIMARY)

    code = '''@TableField("version")
@Version
private Integer version;'''
    add_code_block(slide, right_x, content_top + Inches(0.7), Inches(5.0), Inches(1.5),
                   code, font_size=12)

    lock_explain = [
        "更新时自动 WHERE version = ?",
        "匹配 → 执行操作 + version 自增",
        "不匹配 → 抛出 409 Conflict",
    ]
    add_bullet_list(slide, right_x, content_top + Inches(2.4), Inches(5.0), Inches(1.5),
                   lock_explain, font_size=12)

    # P26 — 测试报告
    page_num += 1
    slide = prs.slides.add_slide(blank_layout)
    content_top = add_title_bar(slide, "测试工作")
    add_page_number(slide, page_num)

    test_cols = [
        ("🔧 功能测试", "5/5 通过", COLOR_PRIMARY,
         ["报销单 CRUD ✅", "行程管理 ✅", "费用合计 ✅", "费用分摊 ✅", "状态流转 ✅"]),
        ("⚡ 性能测试", "3/3 达标", COLOR_SUCCESS,
         ["列表 QPS < 200ms", "详情 QPS < 50ms", "提交 QPS < 500ms"]),
        ("🔒 安全校验", "4/4 通过", COLOR_INNOVATION,
         ["SQL 注入 ✅ 安全", "XSS ✅ 安全", "重复提交 ✅ 限流", "并发冲突 ✅ 409"]),
    ]
    for i, (title, stat, color, items) in enumerate(test_cols):
        x = Inches(0.6 + i * 4.2)
        y = content_top + Inches(0.3)
        card_w = Inches(3.8)
        add_card(slide, x, y, card_w, Inches(4.8), shadow=True)
        # 顶部色条
        add_shape(slide, MSO_SHAPE.RECTANGLE, x, y, card_w, Inches(0.08), fill_color=color)
        # 标题
        add_textbox(slide, x, y + Inches(0.2), card_w, Inches(0.4),
                    title, font_size=16, bold=True, color=color, alignment=PP_ALIGN.CENTER)
        # 通过数
        add_textbox(slide, x, y + Inches(0.6), card_w, Inches(0.4),
                    stat, font_size=20, bold=True, color=COLOR_SUCCESS, alignment=PP_ALIGN.CENTER)
        # 详细项
        add_bullet_list(slide, x + Inches(0.3), y + Inches(1.2), card_w - Inches(0.6),
                       Inches(3.2), items, font_size=12, prefix="▸")

    # P27 — 杨俊杰小结 + 过渡
    page_num += 1
    slide = prs.slides.add_slide(blank_layout)
    add_shape(slide, MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H, fill_color=COLOR_BG)
    add_page_number(slide, page_num)

    add_textbox(slide, Inches(1), Inches(1.2), Inches(11.3), Inches(0.6),
                "杨俊杰 — 模块总结", font_size=26, bold=True, color=COLOR_PRIMARY,
                alignment=PP_ALIGN.CENTER)
    # 分隔线
    add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(4), Inches(1.85), Inches(5.3), Inches(0.03),
              fill_color=COLOR_ACCENT)

    summary_yjj = [
        "核心业务：报销单 CRUD + 行程管理 + 费用合计/分摊 + 状态流转，共 10+ 个接口",
        "技术亮点：分钟级重叠检测 · 首行锁定 · 银行家舍入 · 乐观锁",
        "测试覆盖：功能 5 类 + 性能 3 项 + 安全 4 项",
    ]
    add_bullet_list(slide, Inches(1.5), Inches(2.2), Inches(10), Inches(2.5),
                   summary_yjj, font_size=14, prefix="▸", bold_prefix=True)

    # 过渡
    add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(4), Inches(4.8), Inches(5.3), Inches(0.03),
              fill_color=COLOR_ACCENT)
    add_textbox(slide, Inches(1), Inches(5.0), Inches(11.3), Inches(0.6),
                "接下来 ↓ 吴云龙 — 补助日历 + 数据库设计",
                font_size=18, bold=True, color=COLOR_PRIMARY, alignment=PP_ALIGN.CENTER)

    # ════════════════════════════════════════════════════════
    # Part D — 吴云龙个人专项
    # ════════════════════════════════════════════════════════

    # P28 — 吴云龙个人封面
    page_num += 1
    slide = prs.slides.add_slide(blank_layout)
    add_shape(slide, MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H, fill_color=COLOR_BG)
    add_page_number(slide, page_num)
    add_textbox(slide, Inches(1), Inches(1.5), Inches(11.3), Inches(0.8),
                "吴云龙 个人专项汇报", font_size=32, bold=True, color=COLOR_PRIMARY,
                alignment=PP_ALIGN.CENTER)

    wyl_modules = [
        ("📅", "补助日历", "CRUD + 批量更新 + 联动回写", COLOR_PRIMARY),
        ("🗄️", "数据库设计", "5 张表 + 迁移 V1→V6", COLOR_ACCENT),
    ]
    for i, (icon, title, desc, color) in enumerate(wyl_modules):
        x = Inches(2.0 + i * 5.0)
        y = Inches(3.0)
        add_card(slide, x, y, Inches(4.2), Inches(2.5), shadow=True, border_color=color)
        add_textbox(slide, x, y + Inches(0.3), Inches(4.2), Inches(0.6),
                    icon, font_size=36, alignment=PP_ALIGN.CENTER)
        add_textbox(slide, x, y + Inches(1.0), Inches(4.2), Inches(0.4),
                    title, font_size=18, bold=True, color=color, alignment=PP_ALIGN.CENTER)
        add_textbox(slide, x, y + Inches(1.5), Inches(4.2), Inches(0.4),
                    desc, font_size=12, color=COLOR_TEXT_SEC, alignment=PP_ALIGN.CENTER)

    # P29 — 补助日历业务背景
    page_num += 1
    slide = prs.slides.add_slide(blank_layout)
    content_top = add_title_bar(slide, "补助日历 — 业务背景")
    add_page_number(slide, page_num)

    # 左侧场景
    scenarios = [
        ("6/1 到达北京，晚餐自己解决", "勾选餐补 ✅"),
        ("6/2 全天会议，公司提供餐食", "不勾餐补 ❌"),
        ("6/3 坐高铁回武汉", "勾选交通补 ✅"),
    ]
    add_textbox(slide, Inches(0.6), content_top + Inches(0.3), Inches(5), Inches(0.35),
                "出差场景举例：", font_size=14, bold=True, color=COLOR_PRIMARY)
    for i, (scene, result) in enumerate(scenarios):
        y = content_top + Inches(0.8 + i * 0.9)
        add_card(slide, Inches(0.6), y, Inches(5.5), Inches(0.75), border_color=COLOR_LT2)
        add_textbox(slide, Inches(0.9), y + Inches(0.05), Inches(3.5), Inches(0.35),
                    scene, font_size=12, color=COLOR_TEXT)
        add_textbox(slide, Inches(0.9), y + Inches(0.38), Inches(4.5), Inches(0.3),
                    result, font_size=12, bold=True,
                    color=COLOR_SUCCESS if "✅" in result else COLOR_ERROR)

    # 右侧日历矩阵
    right_x = Inches(6.8)
    add_textbox(slide, right_x, content_top + Inches(0.3), Inches(5), Inches(0.35),
                "日历矩阵示意", font_size=14, bold=True, color=COLOR_PRIMARY)
    # 简单表格模拟
    matrix_data = [
        ["", "6/1", "6/2", "6/3"],
        ["行程A 餐补", "✅", "", "✅"],
        ["行程A 交通补", "✅", "", "✅"],
        ["行程A 通讯补", "✅", "✅", "✅"],
    ]
    add_table(slide, right_x, content_top + Inches(0.8), Inches(5.5), Inches(2.2),
              4, 4, matrix_data, col_widths=[Inches(2.0), Inches(1.1), Inches(1.1), Inches(1.1)])

    add_textbox(slide, right_x, content_top + Inches(3.2), Inches(5.5), Inches(0.3),
                "✅ = 已勾选  |  空 = 未勾选  |  每天 × 每种补助 = 独立选择",
                font_size=10, color=COLOR_TEXT_SEC)

    # P30 — 补助日历数据库设计
    page_num += 1
    slide = prs.slides.add_slide(blank_layout)
    content_top = add_title_bar(slide, "补助日历 — 数据库设计")
    add_page_number(slide, page_num)

    # 左侧表结构
    add_textbox(slide, Inches(0.6), content_top + Inches(0.2), Inches(5), Inches(0.35),
                "reimbursement_allowance（补助主表）", font_size=13, bold=True, color=COLOR_PRIMARY)
    allow_fields = ["id (PK)", "reimbursement_id (FK)", "city_name",
                    "meal_total", "transport_total", "communication_total"]
    add_bullet_list(slide, Inches(0.8), content_top + Inches(0.6), Inches(5), Inches(2.0),
                   allow_fields, font_size=11, prefix="•", bold_prefix=True)

    # 箭头
    add_textbox(slide, Inches(2.5), content_top + Inches(2.8), Inches(1), Inches(0.4),
                "↓ 1:N", font_size=14, bold=True, color=COLOR_ACCENT, alignment=PP_ALIGN.CENTER)

    add_textbox(slide, Inches(0.6), content_top + Inches(3.2), Inches(5), Inches(0.35),
                "reimbursement_allowance_calendar（日历明细）", font_size=13,
                bold=True, color=COLOR_PRIMARY)
    cal_fields = ["id (PK)", "allowance_id (FK)", "calendar_date",
                  "selected", "meal_amount", "transport_amount"]
    add_bullet_list(slide, Inches(0.8), content_top + Inches(3.6), Inches(5), Inches(2.0),
                   cal_fields, font_size=11, prefix="•", bold_prefix=True)

    # 右侧设计要点
    right_x = Inches(7.0)
    add_textbox(slide, right_x, content_top + Inches(0.2), Inches(5), Inches(0.35),
                "设计要点", font_size=15, bold=True, color=COLOR_PRIMARY)
    design_points = [
        "1. 两级关联：报销单 → 补助 → 日历",
        "   形成一对多层级关系",
        "2. 日历粒度：精确到天",
        "   每条记录代表一天的补助选择",
        "3. 金额独立存储：",
        "   每条日历记录存储金额，而非仅 boolean",
    ]
    add_bullet_list(slide, right_x, content_top + Inches(0.7), Inches(5.5), Inches(3.0),
                   design_points, font_size=12, prefix="▸")

    # P31 — 日历接口 + 批量更新
    page_num += 1
    slide = prs.slides.add_slide(blank_layout)
    content_top = add_title_bar(slide, "日历管理接口 + 批量更新设计")
    add_page_number(slide, page_num)

    cal_api_data = [
        ["方法", "路径", "功能"],
        ["GET", ".../calendar", "日历列表"],
        ["POST", ".../calendar", "新增（幂等）"],
        ["PUT", ".../calendar", "批量更新"],
        ["PUT", ".../calendar/{id}", "单条更新"],
        ["DELETE", ".../calendar/{id}", "删除日历项"],
    ]
    add_table(slide, Inches(0.6), content_top + Inches(0.2), Inches(6.0), Inches(2.8),
              6, 3, cal_api_data, col_widths=[Inches(1.2), Inches(2.5), Inches(2.3)])

    # 右侧批量更新 JSON
    right_x = Inches(7.2)
    add_textbox(slide, right_x, content_top + Inches(0.2), Inches(5), Inches(0.35),
                "💡 批量更新请求体", font_size=14, bold=True, color=COLOR_INNOVATION)
    batch_json = '''{
  "calendarItems": [
    {"id": 1, "selected": true,
     "mealAmount": 100.00},
    {"id": 2, "selected": false,
     "mealAmount": 0}
  ]
}'''
    add_code_block(slide, right_x, content_top + Inches(0.7), Inches(5.2), Inches(2.5),
                   batch_json, font_size=11)
    add_textbox(slide, right_x, content_top + Inches(3.3), Inches(5.2), Inches(0.3),
                "→ 一次请求更新所有变更，事务保证一致性",
                font_size=11, color=COLOR_TEXT_SEC)

    # P32 — 幂等性设计
    page_num += 1
    slide = prs.slides.add_slide(blank_layout)
    content_top = add_title_bar(slide, "幂等性设计")
    add_innovation_badge(slide, Inches(4.0), content_top - Inches(0.55))
    add_page_number(slide, page_num)

    # 流程图
    add_box_with_text(slide, Inches(1.5), content_top + Inches(0.5),
                     Inches(3.0), Inches(0.7),
                     "POST 新增日历项", bg_color=COLOR_PRIMARY, font_size=13)
    add_textbox(slide, Inches(2.7), content_top + Inches(1.3), Inches(0.5), Inches(0.4),
                "↓", font_size=18, color=COLOR_TEXT_SEC)
    add_box_with_text(slide, Inches(0.5), content_top + Inches(1.8),
                     Inches(5.5), Inches(0.7),
                     "同一天 + 同一补助类型是否已存在？", bg_color=COLOR_ACCENT, font_size=13)

    # 分支 — 存在
    add_textbox(slide, Inches(0.8), content_top + Inches(2.7), Inches(1.5), Inches(0.3),
                "存在 ↓", font_size=12, bold=True, color=COLOR_ERROR)
    add_card(slide, Inches(0.3), content_top + Inches(3.0), Inches(3.0), Inches(0.7),
            fill=RGBColor(0xFD, 0xF2, 0xF2), border_color=COLOR_ERROR)
    add_textbox(slide, Inches(0.5), content_top + Inches(3.1), Inches(2.6), Inches(0.5),
                "返回已有记录\n（不重复创建）", font_size=11, color=COLOR_ERROR,
                alignment=PP_ALIGN.CENTER)

    # 分支 — 不存在
    add_textbox(slide, Inches(4.0), content_top + Inches(2.7), Inches(2.0), Inches(0.3),
                "不存在 ↓", font_size=12, bold=True, color=COLOR_SUCCESS)
    add_card(slide, Inches(3.5), content_top + Inches(3.0), Inches(3.0), Inches(0.7),
            fill=RGBColor(0xE8, 0xF8, 0xEF), border_color=COLOR_SUCCESS)
    add_textbox(slide, Inches(3.7), content_top + Inches(3.1), Inches(2.6), Inches(0.5),
                "创建新记录", font_size=11, color=COLOR_SUCCESS,
                alignment=PP_ALIGN.CENTER)

    # 右侧说明
    right_x = Inches(7.2)
    add_textbox(slide, right_x, content_top + Inches(0.5), Inches(5.5), Inches(0.4),
                "幂等性 = 无论调用多少次，结果一致",
                font_size=16, bold=True, color=COLOR_PRIMARY)
    add_textbox(slide, right_x, content_top + Inches(1.2), Inches(5.5), Inches(1.5),
                "通过「先查询后创建」策略，\n"
                "保证重复请求不会产生重复数据。\n\n"
                "即使网络抖动导致前端重复发送，\n"
                "后端也能正确处理。",
                font_size=12, color=COLOR_TEXT_SEC)

    # P33 — 联动回写机制
    page_num += 1
    slide = prs.slides.add_slide(blank_layout)
    content_top = add_title_bar(slide, "补助金额联动回写")
    add_page_number(slide, page_num)

    # 横向三步流程
    flow_steps = [
        ("日历选择变化", "遍历所有日历记录", COLOR_PRIMARY),
        ("补助汇总重算", "统计 selected=true 的金额", COLOR_ACCENT),
        ("报销总额更新", "前端重新调用费用合计接口", COLOR_SUCCESS),
    ]
    for i, (title, desc, color) in enumerate(flow_steps):
        x = Inches(0.6 + i * 4.3)
        y = content_top + Inches(0.4)
        add_box_with_text(slide, x, y, Inches(3.5), Inches(0.6),
                         title, bg_color=color, font_size=13)
        add_textbox(slide, x, y + Inches(0.7), Inches(3.5), Inches(0.4),
                    desc, font_size=10, color=COLOR_TEXT_SEC, alignment=PP_ALIGN.CENTER)
        if i < 2:
            add_flow_arrow(slide, x + Inches(3.6), y + Inches(0.15),
                          Inches(0.5), Inches(0.25))

    # 数据一致性保证
    add_textbox(slide, Inches(0.6), content_top + Inches(2.0), Inches(10), Inches(0.4),
                "数据一致性保证", font_size=16, bold=True, color=COLOR_PRIMARY)
    guarantees = [
        ("数据库事务", "同一事务中完成全部操作", COLOR_PRIMARY),
        ("乐观锁", "version 字段防并发冲突", COLOR_ACCENT),
        ("银行家舍入", "精度一致", COLOR_SUCCESS),
    ]
    for i, (title, desc, color) in enumerate(guarantees):
        x = Inches(0.6 + i * 4.3)
        y = content_top + Inches(2.5)
        add_card(slide, x, y, Inches(3.5), Inches(1.8), shadow=True, border_color=color)
        add_textbox(slide, x, y + Inches(0.2), Inches(3.5), Inches(0.4),
                    title, font_size=15, bold=True, color=color, alignment=PP_ALIGN.CENTER)
        add_textbox(slide, x + Inches(0.3), y + Inches(0.7), Inches(2.9), Inches(0.8),
                    desc, font_size=12, color=COLOR_TEXT_SEC, alignment=PP_ALIGN.CENTER)

    # P34 — 数据库设计原则 + 5 张表总览
    page_num += 1
    slide = prs.slides.add_slide(blank_layout)
    content_top = add_title_bar(slide, "报销模块数据库设计")
    add_page_number(slide, page_num)

    # 设计原则标签
    principles = ["主子表分离", "外键关联", "冗余字段", "版本控制"]
    for i, p in enumerate(principles):
        x = Inches(0.6 + i * 3.1)
        add_box_with_text(slide, x, content_top + Inches(0.3), Inches(2.6), Inches(0.5),
                         p, bg_color=COLOR_ACCENT, font_size=12)

    # 5 张表表格
    table_data_db = [
        ["表名", "作用", "关联方式"],
        ["reimbursement", "报销单主表", "—"],
        ["reimbursement_travel_record", "行程子表", "FK → reimbursement.id"],
        ["reimbursement_allowance", "补助子表", "FK → reimbursement.id"],
        ["reimbursement_allowance_calendar", "日历明细", "FK → reimbursement_allowance.id"],
        ["reimbursement_cost_allocation", "分摊子表", "FK → reimbursement.id"],
    ]
    add_table(slide, Inches(0.6), content_top + Inches(1.1), Inches(12), Inches(3.2),
              6, 3, table_data_db, col_widths=[Inches(5.0), Inches(2.5), Inches(4.5)])

    # P35 — ER 图
    page_num += 1
    slide = prs.slides.add_slide(blank_layout)
    content_top = add_title_bar(slide, "报销模块 ER 关系图")
    add_page_number(slide, page_num)

    # 中心表
    center_x = Inches(4.5)
    center_y = content_top + Inches(1.5)
    add_card(slide, center_x, center_y, Inches(3.5), Inches(2.5),
            border_color=COLOR_PRIMARY, shadow=True)
    add_textbox(slide, center_x, center_y + Inches(0.1), Inches(3.5), Inches(0.35),
                "reimbursement", font_size=14, bold=True, color=COLOR_PRIMARY,
                alignment=PP_ALIGN.CENTER)
    fields_main = ["id (PK)", "content (JSON)", "version", "status", "reimburser_name", "..."]
    add_bullet_list(slide, center_x + Inches(0.3), center_y + Inches(0.5),
                   Inches(2.9), Inches(1.8), fields_main, font_size=10, prefix="•")

    # 四个子表
    child_tables = [
        ("travel_record", Inches(0.5), Inches(1.2), COLOR_ACCENT),
        ("allowance", Inches(0.5), Inches(4.0), COLOR_SUCCESS),
        ("allowance_calendar", Inches(9.5), Inches(4.0), COLOR_INNOVATION),
        ("cost_allocation", Inches(9.5), Inches(1.2), COLOR_ERROR),
    ]
    for name, x, y, color in child_tables:
        add_card(slide, x, y, Inches(3.0), Inches(1.2), border_color=color)
        add_textbox(slide, x, y + Inches(0.1), Inches(3.0), Inches(0.3),
                    f"reimbursement_{name}", font_size=10, bold=True, color=color,
                    alignment=PP_ALIGN.CENTER)
        add_textbox(slide, x + Inches(0.2), y + Inches(0.45), Inches(2.6), Inches(0.6),
                    "FK → reimbursement.id", font_size=9, color=COLOR_TEXT_SEC,
                    alignment=PP_ALIGN.CENTER)

    # 关系标注
    add_textbox(slide, Inches(3.5), Inches(2.5), Inches(1.0), Inches(0.3),
                "1:N →", font_size=11, color=COLOR_TEXT_SEC, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(3.5), Inches(4.5), Inches(1.0), Inches(0.3),
                "1:N →", font_size=11, color=COLOR_TEXT_SEC, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(8.0), Inches(2.5), Inches(1.0), Inches(0.3),
                "← 1:N", font_size=11, color=COLOR_TEXT_SEC, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(8.0), Inches(4.5), Inches(1.0), Inches(0.3),
                "← 1:N", font_size=11, color=COLOR_TEXT_SEC, alignment=PP_ALIGN.CENTER)

    add_textbox(slide, Inches(0.6), Inches(6.6), Inches(10), Inches(0.3),
                "来源：docs/database-schema-mermaid.md", font_size=10, color=COLOR_TEXT_SEC)

    # P36 — 数据库版本化迁移
    page_num += 1
    slide = prs.slides.add_slide(blank_layout)
    content_top = add_title_bar(slide, "数据库版本化迁移 V1 → V6")
    add_page_number(slide, page_num)

    # 左侧时间轴
    versions = [
        ("V1", "初始化"),
        ("V2", "报销模块 5 张表创建"),
        ("V3", "样本数据"),
        ("V4", "批量性能测试数据"),
        ("V5", "行程表增加 datetime 字段"),
        ("V6", "回填已有数据的时间字段"),
    ]
    for i, (ver, desc) in enumerate(versions):
        y = content_top + Inches(0.3 + i * 0.7)
        add_status_circle(slide, Inches(0.8), y, Inches(0.25), COLOR_PRIMARY)
        add_textbox(slide, Inches(1.2), y - Inches(0.05), Inches(1.0), Inches(0.3),
                    ver, font_size=13, bold=True, color=COLOR_PRIMARY)
        add_textbox(slide, Inches(2.2), y - Inches(0.05), Inches(4.0), Inches(0.3),
                    desc, font_size=12, color=COLOR_TEXT)
        # 竖线
        if i < 5:
            add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0.9), y + Inches(0.25),
                     Inches(0.02), Inches(0.45), fill_color=COLOR_LT2)

    # 右侧迁移策略
    right_x = Inches(7.0)
    add_textbox(slide, right_x, content_top + Inches(0.2), Inches(5), Inches(0.35),
                "平滑迁移策略（以 V5→V6 为例）", font_size=14, bold=True, color=COLOR_PRIMARY)
    migration_steps = [
        "① ALTER TABLE ADD COLUMN\n   （允许 NULL）",
        "② UPDATE\n   回填已有数据",
        "③ ALTER NOT NULL\n   （可选）",
    ]
    add_bullet_list(slide, right_x, content_top + Inches(0.7), Inches(5.5), Inches(2.5),
                   migration_steps, font_size=12, prefix="", bold_prefix=True)

    add_green_check_items(slide, right_x, content_top + Inches(3.3), Inches(5.5),
                         ["先加字段 → 再回填 → 最后验证", "不是一步到位改表"],
                         font_size=12)

    # P37 — 吴云龙测试报告
    page_num += 1
    slide = prs.slides.add_slide(blank_layout)
    content_top = add_title_bar(slide, "测试工作")
    add_page_number(slide, page_num)

    test_cols_wyl = [
        ("🔧 功能测试", "5/5 通过", COLOR_PRIMARY,
         ["日历 CRUD ✅", "幂等性 ✅", "联动回写 ✅", "金额精度 ✅", "DB 迁移 ✅"]),
        ("⚡ 性能测试", "2/2 达标", COLOR_SUCCESS,
         ["批量更新 < 100ms", "补助生成 < 200ms"]),
        ("🔒 安全校验", "3/3 通过", COLOR_INNOVATION,
         ["SQL 注入 ✅ 安全", "越权访问 ✅ 校验", "事务回滚 ✅ 整体回滚"]),
    ]
    for i, (title, stat, color, items) in enumerate(test_cols_wyl):
        x = Inches(0.6 + i * 4.2)
        y = content_top + Inches(0.3)
        card_w = Inches(3.8)
        add_card(slide, x, y, card_w, Inches(4.8), shadow=True)
        add_shape(slide, MSO_SHAPE.RECTANGLE, x, y, card_w, Inches(0.08), fill_color=color)
        add_textbox(slide, x, y + Inches(0.2), card_w, Inches(0.4),
                    title, font_size=16, bold=True, color=color, alignment=PP_ALIGN.CENTER)
        add_textbox(slide, x, y + Inches(0.6), card_w, Inches(0.4),
                    stat, font_size=20, bold=True, color=COLOR_SUCCESS, alignment=PP_ALIGN.CENTER)
        add_bullet_list(slide, x + Inches(0.3), y + Inches(1.2), card_w - Inches(0.6),
                       Inches(3.2), items, font_size=12, prefix="▸")

    # P38 — 结尾致谢页
    page_num += 1
    slide = prs.slides.add_slide(blank_layout)
    # 背景
    add_shape(slide, MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H, fill_color=COLOR_BG)

    # 致谢文字
    add_textbox(slide, Inches(1), Inches(1.2), Inches(11.3), Inches(0.8),
                "感谢各位老师！", font_size=36, bold=True, color=COLOR_PRIMARY,
                alignment=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(1), Inches(2.1), Inches(11.3), Inches(0.6),
                "Thank You!", font_size=22, color=COLOR_TEXT_SEC, alignment=PP_ALIGN.CENTER)

    # 信息卡片
    card_x = Inches(2.5)
    card_y = Inches(3.0)
    card_w = Inches(8.3)
    card_h = Inches(2.8)
    add_card(slide, card_x, card_y, card_w, card_h, shadow=True, border_color=COLOR_ACCENT)

    info_lines = [
        ("TripFlow — 企业差旅报销流程管理平台", 16, True, COLOR_PRIMARY),
        ("", 8, False, COLOR_TEXT),
        ("吴汉东 · 杨俊杰 · 吴云龙", 14, False, COLOR_TEXT),
        ("", 8, False, COLOR_TEXT),
        ("技术栈：Spring Boot + Vue 3 + MySQL + Redis", 13, False, COLOR_TEXT_SEC),
        ("关键数据：30+ API · 11 张表 · 4000+ 行前端代码", 13, False, COLOR_TEXT_SEC),
    ]
    y_offset = card_y + Inches(0.3)
    for text, fs, bold, color in info_lines:
        if text:
            add_textbox(slide, card_x, y_offset, card_w, Inches(0.4),
                        text, font_size=fs, bold=bold, color=color, alignment=PP_ALIGN.CENTER)
        y_offset += Inches(0.38)

    # 欢迎提问
    add_textbox(slide, Inches(1), Inches(6.2), Inches(11.3), Inches(0.6),
                "欢迎提问 🙋", font_size=20, color=COLOR_TEXT_SEC, alignment=PP_ALIGN.CENTER)

    # ════════════════════════════════════════════════════════
    # 保存
    # ════════════════════════════════════════════════════════
    os.makedirs(OUTPUT_DIR, exist_ok=True)
    prs.save(OUTPUT_PATH)
    print(f"✅ PPT 已生成：{OUTPUT_PATH}")
    print(f"   共 {page_num} 页")


if __name__ == "__main__":
    generate_presentation()
